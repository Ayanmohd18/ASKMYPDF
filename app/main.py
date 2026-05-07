import streamlit as st
import os
import tempfile
import datetime
from pathlib import Path

# MUST BE THE FIRST STREAMLIT COMMAND
st.set_page_config(
    page_title="AskMyPDF",
    page_icon="📄",
    layout="wide",
    initial_sidebar_state="expanded"
)

# --- Premium Custom CSS ---
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;600;800&display=swap');

/* Global Styles & Typography */
html, body, [class*="css"], .stApp {
    font-family: 'Outfit', sans-serif !important;
}

/* Vibrant & Modern Accent Gradients */
:root {
    --accent-gradient: linear-gradient(135deg, #6366f1 0%, #a855f7 100%);
    --glass-bg: rgba(255, 255, 255, 0.7);
    --glass-border: rgba(255, 255, 255, 0.4);
    --shadow-color: rgba(99, 102, 241, 0.15);
}

@media (prefers-color-scheme: dark) {
    :root {
        --glass-bg: rgba(15, 23, 42, 0.6);
        --glass-border: rgba(255, 255, 255, 0.08);
        --shadow-color: rgba(0, 0, 0, 0.3);
    }
}

/* Dynamic Ambient Background Glow */
.stApp {
    background: radial-gradient(circle at 90% 10%, rgba(99, 102, 241, 0.1), transparent 40%),
                radial-gradient(circle at 10% 90%, rgba(168, 85, 247, 0.08), transparent 40%);
}

/* Smooth State Transitions */
* {
    transition: all 0.3s ease-in-out;
}

/* Glassmorphic Sidebar */
[data-testid="stSidebar"] {
    background: var(--glass-bg) !important;
    backdrop-filter: blur(12px) !important;
    border-right: 1px solid var(--glass-border) !important;
}

/* Uber-Aesthetic Premium Buttons */
.stButton > button {
    background: var(--accent-gradient) !important;
    color: white !important;
    border: none !important;
    border-radius: 12px !important;
    font-weight: 600 !important;
    letter-spacing: 0.5px !important;
    padding: 12px 24px !important;
    box-shadow: 0 4px 15px var(--shadow-color) !important;
}

.stButton > button:hover {
    transform: translateY(-3px) scale(1.02) !important;
    box-shadow: 0 8px 25px rgba(99, 102, 241, 0.4) !important;
}

/* Input Chat Box Glow */
.stChatInputContainer {
    border-radius: 24px !important;
    background: var(--glass-bg) !important;
    backdrop-filter: blur(12px) !important;
    border: 1px solid var(--glass-border) !important;
    box-shadow: 0 10px 30px var(--shadow-color) !important;
}

.stChatInputContainer:focus-within {
    border-color: #6366f1 !important;
    box-shadow: 0 10px 35px rgba(99, 102, 241, 0.25) !important;
}

/* Chat Bubbles */
.stChatMessage {
    border-radius: 18px !important;
    border: 1px solid var(--glass-border) !important;
    padding: 16px !important;
    background: var(--glass-bg) !important;
    backdrop-filter: blur(10px) !important;
    margin-bottom: 16px !important;
    box-shadow: 0 4px 15px var(--shadow-color) !important;
}

/* Expanders (Sources) */
.streamlit-expanderHeader {
    background: var(--glass-bg) !important;
    border-radius: 12px !important;
    border: 1px solid var(--glass-border) !important;
    font-weight: 600 !important;
}

/* Smooth Micro-Animations for Flashcards */
.flashcard {
    perspective: 1000px;
    margin-bottom: 20px;
}
.flashcard-inner {
    position: relative;
    width: 100%;
    height: 200px;
    text-align: center;
    transition: transform 0.8s cubic-bezier(0.175, 0.885, 0.32, 1.275);
    transform-style: preserve-3d;
    cursor: pointer;
}
.flashcard:hover .flashcard-inner {
    transform: rotateY(180deg);
}
.flashcard-front, .flashcard-back {
    position: absolute;
    width: 100%;
    height: 100%;
    backface-visibility: hidden;
    display: flex;
    align-items: center;
    justify-content: center;
    padding: 20px;
    border-radius: 20px;
    border: 1px solid var(--glass-border);
    box-shadow: 0 10px 25px var(--shadow-color);
    backdrop-filter: blur(12px);
}
.flashcard-front {
    background: var(--glass-bg);
    font-weight: 600;
    font-size: 1.3em;
    color: inherit;
}
.flashcard-back {
    background: var(--accent-gradient);
    color: white;
    transform: rotateY(180deg);
    font-size: 1.1em;
}
</style>

""", unsafe_allow_html=True)


from app.monitor import get_stats, get_tracker
from app.memory import ConversationMemory
from app.retriever import retrieve, RetrievedChunk
from app.generator import generate
from app import generator
from app import ingestion
from app import vector_store
from dotenv import load_dotenv

load_dotenv()

# Initialize vector store
index_dir = Path(os.getenv("INDEX_DIR", "data/indexes"))
index_dir.mkdir(parents=True, exist_ok=True)
vector_store.initialize(index_dir)

# Initialize Session State
if "messages" not in st.session_state:
    st.session_state.messages = []
if "memory" not in st.session_state:
    st.session_state.memory = ConversationMemory()
if "backend" not in st.session_state:
    st.session_state.backend = os.getenv("LLM_BACKEND", "gemini")
if "use_hyde" not in st.session_state:
    st.session_state.use_hyde = False
if "use_decomposition" not in st.session_state:
    st.session_state.use_decomposition = False
if "last_query_latency" not in st.session_state:
    st.session_state.last_query_latency = 0.0
if "ingested_docs" not in st.session_state:
    st.session_state.ingested_docs = {}
    # Populate initial from vector store if available
    for doc in vector_store.get_all_doc_names():
        st.session_state.ingested_docs[doc] = {
            "chunks": "Loaded from disk",
            "timestamp": "Pre-existing"
        }

def render_source_cards(results: list[RetrievedChunk]):
    st.markdown("**📎 Sources:**")
    for i, rc in enumerate(results, start=1):
        with st.expander(
            f"[{i}] {rc.chunk.doc_name} — Page {rc.chunk.page_number} "
            f"(score: {rc.reranker_score:.3f})"
        ):
            st.markdown(f"**Document:** {rc.chunk.doc_name}")
            st.markdown(f"**Page:** {rc.chunk.page_number}")
            col1, col2, col3 = st.columns(3)
            col1.metric("Reranker", f"{rc.reranker_score:.3f}")
            col2.metric("FAISS", f"{rc.faiss_score:.3f}")
            col3.metric("BM25", f"{rc.bm25_score:.3f}")
            st.divider()
            st.caption("Relevant passage:")
            st.markdown(f"> {rc.chunk.text[:500]}...")

# ─── SIDEBAR ────────────────────────────────────────────
with st.sidebar:
    st.header("📁 Document Management")
    uploaded_files = st.file_uploader("Upload Files", type=["pdf", "docx", "txt"], accept_multiple_files=True, key="uploaded_files")
    url_input = st.text_input("Ingest Website URL")
    
    if st.button("⚙️ Ingest Documents"):
        if not uploaded_files and not url_input.strip():
            st.warning("Please upload files or provide a URL.")
        else:
            progress_bar = st.progress(0)
            total_items = len(uploaded_files) + (1 if url_input.strip() else 0)
            all_chunks = []
            
            # Process uploaded files
            for idx, file in enumerate(uploaded_files):
                try:
                    with tempfile.TemporaryDirectory() as tmpdir:
                        static_dir = Path("static")
                        static_dir.mkdir(exist_ok=True)
                        file_bytes = file.read()
                        
                        tmp_path = Path(tmpdir) / file.name
                        tmp_path.write_bytes(file_bytes)
                        
                        static_path = static_dir / file.name
                        static_path.write_bytes(file_bytes)
                        
                        chunks = ingestion.ingest_file(tmp_path)
                        all_chunks.extend(chunks)
                    
                    st.session_state.ingested_docs[file.name] = {
                        "chunks": len(chunks),
                        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    }
                except Exception as e:
                    st.error(f"Error processing {file.name}: {e}")
                progress_bar.progress((idx + 1) / total_items)
            
            # Process URL
            if url_input.strip():
                try:
                    chunks = ingestion.ingest_url(url_input.strip())
                    all_chunks.extend(chunks)
                    st.session_state.ingested_docs[url_input.strip()] = {
                        "chunks": len(chunks),
                        "timestamp": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    }
                except Exception as e:
                    st.error(f"Error processing URL: {e}")
                progress_bar.progress(1.0)
                
            if all_chunks:
                try:
                    vector_store.add_chunks(all_chunks)
                    st.success(f"Successfully indexed {len(all_chunks)} chunks!")
                except Exception as e:
                    st.error(f"Error updating vector store: {e}")
                    
    st.header("📚 Indexed Documents")
    if not st.session_state.ingested_docs:
        st.info("No documents indexed yet.")
    else:
        df_data = []
        for doc, info in st.session_state.ingested_docs.items():
            df_data.append({
                "Document": doc,
                "Chunks": info["chunks"],
                "Ingested At": info["timestamp"]
            })
        st.dataframe(df_data, use_container_width=True)
        
    if st.button("🗑️ Clear All Documents"):
        vector_store.clear()
        st.session_state.ingested_docs = {}
        st.success("Vector store cleared.")
        
    st.header("🤖 LLM Backend")
    st.session_state.backend = st.radio(
        "Select Backend:",
        ["gemini", "hf", "ollama"],
        index=["gemini", "hf", "ollama"].index(st.session_state.backend) if st.session_state.backend in ["gemini", "hf", "ollama"] else 0
    )
    
    st.header("🧠 Advanced Retrieval")
    st.session_state.use_hyde = st.toggle(
        "Enable HyDE (Hypothetical Embeddings)", 
        value=st.session_state.use_hyde, 
        help="Improves semantic search by asking the LLM to hallucinate a plausible answer first, and then searching the vector DB for it."
    )
    st.session_state.use_decomposition = st.toggle(
        "Enable Query Decomposition", 
        value=st.session_state.use_decomposition, 
        help="Breaks complex questions into multiple sub-queries to retrieve a broader set of relevant context."
    )
    
    st.header("🎤 Voice Interface")
    if hasattr(st, "audio_input"):
        voice_query_audio = st.audio_input("Record a voice query")
        if voice_query_audio is not None:
            if st.session_state.get("last_voice_audio") != voice_query_audio:
                with st.spinner("Transcribing voice via Whisper..."):
                    import whisper
                    import tempfile
                    with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
                        f.write(voice_query_audio.getvalue())
                        tmp_audio_path = f.name
                    
                    # Load Whisper model (base for speed)
                    model = whisper.load_model("base")
                    result = model.transcribe(tmp_audio_path)
                    st.session_state.voice_query_text = result["text"].strip()
                    st.session_state.last_voice_audio = voice_query_audio
                    
                    import os
                    try:
                        os.remove(tmp_audio_path)
                    except Exception:
                        pass
    else:
        st.info("Voice Input is not supported on this Streamlit version (upgrade required).")
    
    st.header("📊 System Monitor")
    stats = get_stats()
    col1, col2, col3 = st.columns(3)
    col1.metric("CPU", f"{stats.cpu_percent:.1f}%")
    col2.metric("RAM", f"{stats.ram_used_mb:.0f} MB")
    col3.metric("Latency", f"{st.session_state.last_query_latency:.2f}s")
    st.caption(f"RAM: {stats.ram_used_mb:.0f}/{stats.ram_total_mb:.0f} MB total")

# ─── MAIN AREA ──────────────────────────────────────────
col_chat, col_studio = st.columns([5, 3])

with col_chat:
    st.title("📄 AskMyPDF Chat")
    st.caption("Chat with your multi-modal documents.")

    col_a, col_b = st.columns([1, 4])
    with col_a:
        if st.button("🧹 Clear Chat"):
            st.session_state.messages = []
            st.session_state.memory.clear()
    with col_b:
        st.caption(f"Documents: {len(st.session_state.ingested_docs)} | Chunks: {vector_store.get_chunk_count()}")

    for msg in st.session_state.messages:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])
            if msg.get("audio"):
                st.audio(msg["audio"], format="audio/mp3")
            if msg["role"] == "assistant" and msg.get("sources"):
                render_source_cards(msg["sources"])

    query = st.chat_input("Ask a question about your docs...")
    
    if st.session_state.get("voice_query_text"):
        query = st.session_state.voice_query_text
        st.session_state.voice_query_text = None

    if query:
        st.session_state.messages.append({"role": "user", "content": query, "sources": []})
        
        with st.chat_message("user"):
            st.markdown(query)
            
        tracker = get_tracker()
        tracker.start()
        
        with st.chat_message("assistant"):
            try:
                reranker_top_k = int(os.getenv("RERANKER_TOP_K", 5))
                
                with st.status("Processing query...", expanded=True) as status:
                    if st.session_state.get("use_decomposition"):
                        status.write("Decomposing complex query...")
                        sub_queries = generator.decompose_query(query, st.session_state.backend)
                        if len(sub_queries) > 1:
                            status.write(f"Generated {len(sub_queries)} sub-queries.")
                            for sq in sub_queries:
                                status.write(f"- {sq}")
                        else:
                            status.write("Query is straightforward, using as-is.")
                            
                        all_results = []
                        seen_chunk_ids = set()
                        
                        for sq in sub_queries:
                            status.write(f"Searching for: *{sq}*")
                            sq_results = retrieve(sq, reranker_top_k=reranker_top_k, use_hyde=st.session_state.use_hyde)
                            for rc in sq_results:
                                if rc.chunk.chunk_id not in seen_chunk_ids:
                                    seen_chunk_ids.add(rc.chunk.chunk_id)
                                    all_results.append(rc)
                        
                        all_results.sort(key=lambda x: x.reranker_score, reverse=True)
                        results = all_results[:10]
                    else:
                        status.write("Searching documents...")
                        results = retrieve(query, reranker_top_k=reranker_top_k, use_hyde=st.session_state.use_hyde)
                        
                    status.write("Generating final answer...")
                    answer = generate(query, results, st.session_state.memory, backend=st.session_state.backend)
                    status.update(label="Response Generated", state="complete", expanded=False)
                    
                st.session_state.memory.add_user(query)
                st.session_state.memory.add_assistant(answer)
                
                # Generate TTS audio using pyttsx3 (local offline TTS)
                audio_bytes = None
                try:
                    import pyttsx3
                    import tempfile
                    import os
                    clean_text = answer.replace("*", "").replace("#", "").replace("_", "")
                    
                    engine = pyttsx3.init()
                    engine.setProperty('rate', 150) # Standard talking speed
                    
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as tmp_file:
                        tmp_path = tmp_file.name
                        
                    engine.save_to_file(clean_text, tmp_path)
                    engine.runAndWait()
                    
                    with open(tmp_path, "rb") as f:
                        audio_bytes = f.read()
                        
                    try:
                        os.remove(tmp_path)
                    except Exception:
                        pass
                except Exception as e:
                    print(f"pyttsx3 TTS Error: {e}")
                
                latency = tracker.stop()
                st.session_state.last_query_latency = latency
                
                st.markdown(answer)
                if audio_bytes:
                    st.audio(audio_bytes, format="audio/mp3")
                if results:
                    render_source_cards(results)
                    
                st.session_state.messages.append({
                    "role": "assistant", 
                    "content": answer, 
                    "sources": results,
                    "audio": audio_bytes
                })
                
                st.rerun()
                    
            except ValueError as e:
                st.warning(str(e))
                tracker.stop()
            except Exception as e:
                st.error(f"Error: {e}")
                tracker.stop()

with col_studio:
    st.title("🎨 Studio")
    st.caption("AI Generated Assets")
    
    tab_audio, tab_slide, tab_video, tab_mindmap, tab_report, tab_flashcard, tab_quiz, tab_info, tab_data = st.tabs([
        "Audio", "Slides", "Video", "Mind Map", "Reports", "Flashcards", "Quiz", "Infographics", "Data"
    ])
    
    with tab_audio:
        st.subheader("Podcast Generation")
        col1, col2 = st.columns(2)
        podcast_length = None
        if col1.button("🎙️ Crispy Response (3-10 mins)"):
            podcast_length = "Crispy Response"
        if col2.button("🎙️ In-Depth Response (5-15 mins)"):
            podcast_length = "In-Depth Response"
            
        if podcast_length:
            if not st.session_state.ingested_docs:
                st.warning("Please index documents first.")
            else:
                with st.spinner(f"Generating {podcast_length} podcast script..."):
                    try:
                        results = retrieve("Extract main topics for a comprehensive podcast.", reranker_top_k=20)
                        script_json = generator.generate_podcast_script(results, st.session_state.backend, podcast_length)
                        
                        st.info("Script generated. Synthesizing audio via ElevenLabs...")
                        audio_bytes = generator.synthesize_audio(script_json)
                        st.session_state["podcast_audio"] = audio_bytes
                        st.success("Podcast ready!")
                    except Exception as e:
                        st.error(f"Error: {e}")
                        
        if "podcast_audio" in st.session_state:
            st.audio(st.session_state["podcast_audio"], format="audio/mp3")
    with tab_slide:
        st.subheader("Presentation Generator")
        num_slides = st.radio("Number of Slides:", [10, 15, 20], horizontal=True)
        template = st.selectbox("Aesthetic Template:", ["Corporate", "Creative", "Minimalist", "Dark Mode"])
        color_palette = st.color_picker("Primary Accent Color", "#0052cc")
        
        if st.button("Generate Slide Deck"):
            if not st.session_state.ingested_docs:
                st.warning("Please index documents first.")
            else:
                with st.spinner("Designing presentation..."):
                    try:
                        results = retrieve("Extract main topics for a comprehensive presentation.", reranker_top_k=20)
                        import json
                        slides_json_str = generator.generate_slides_json(results, st.session_state.backend, num_slides)
                        slides_data = json.loads(slides_json_str.replace('```json', '').replace('```', '').strip())
                        
                        from pptx import Presentation
                        from pptx.util import Inches, Pt
                        from pptx.dml.color import RGBColor
                        
                        prs = Presentation()
                        
                        hex_color = color_palette.lstrip('#')
                        rgb_accent = RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
                        
                        for slide_data in slides_data:
                            slide_layout = prs.slide_layouts[1]
                            slide = prs.slides.add_slide(slide_layout)
                            
                            if template == "Dark Mode":
                                background = slide.background
                                fill = background.fill
                                fill.solid()
                                fill.fore_color.rgb = RGBColor(30, 30, 30)
                                
                            title_shape = slide.shapes.title
                            body_shape = slide.placeholders[1]
                            
                            title_shape.text = slide_data.get("title", "")
                            for p in title_shape.text_frame.paragraphs:
                                p.font.color.rgb = rgb_accent
                                if template == "Dark Mode":
                                    p.font.color.rgb = RGBColor(255, 255, 255)
                                    
                            tf = body_shape.text_frame
                            tf.text = slide_data.get("subtitle", "")
                            
                            for bullet in slide_data.get("bullets", []):
                                p = tf.add_paragraph()
                                p.text = bullet
                                p.level = 1
                                if template == "Dark Mode":
                                    p.font.color.rgb = RGBColor(200, 200, 200)
                                    
                            vs = slide_data.get("visual_suggestion", "")
                            if vs:
                                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
                                tf2 = txBox.text_frame
                                p2 = tf2.add_paragraph()
                                p2.text = f"[Visual: {vs}]"
                                p2.font.size = Pt(12)
                                p2.font.italic = True
                                p2.font.color.rgb = RGBColor(128, 128, 128)
                        
                        import io
                        pptx_io = io.BytesIO()
                        prs.save(pptx_io)
                        st.session_state["pptx_io"] = pptx_io.getvalue()
                        st.success("Slide deck generated!")
                    except Exception as e:
                        st.error(f"Error generating slides: {e}")
                        
        if "pptx_io" in st.session_state:
            st.download_button("📥 Download Presentation (.pptx)", data=st.session_state["pptx_io"], file_name="AskMyPDF_Deck.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    with tab_video:
        st.subheader("Animated Video Overview")
        avatar_style = st.selectbox("Select Avatar:", ["Daisy-in-Tshirt", "Tyler-in-Suit", "Anna_public_3_20240108"])
        if st.button("🎬 Generate Animated Video Overview"):
            if not st.session_state.ingested_docs:
                st.warning("Please index documents first.")
            else:
                with st.spinner("Writing video script..."):
                    try:
                        results = retrieve("Extract main topics for a 60-second video overview.", reranker_top_k=20)
                        script_text = generator.generate_video_script(results, st.session_state.backend)
                        
                        st.info("Sending to Video API (HeyGen). This may take a few minutes...")
                        video_url = generator.generate_video(script_text, avatar_style)
                        st.session_state["video_url"] = video_url
                        st.success("Video generated!")
                    except Exception as e:
                        st.error(f"Error: {e}")
                        
        if "video_url" in st.session_state:
            st.video(st.session_state["video_url"])
        
    with tab_mindmap:
        st.subheader("Interactive Knowledge Graph")
        if st.button("Generate Mind Map"):
            if not st.session_state.ingested_docs:
                st.warning("Please index documents first.")
            else:
                with st.spinner("Extracting graph relationships..."):
                    try:
                        results = retrieve("Extract main entities and their relationships.", reranker_top_k=20)
                        import json
                        map_json_str = generator.generate_mindmap_json(results, st.session_state.backend)
                        map_data = json.loads(map_json_str.replace('```json', '').replace('```', '').strip())
                        st.session_state["mindmap_data"] = map_data
                    except Exception as e:
                        st.error(f"Error extracting graph: {e}")
                        
        if "mindmap_data" in st.session_state:
            try:
                from streamlit_agraph import agraph, Node, Edge, Config
                map_data = st.session_state["mindmap_data"]
                nodes = []
                edges = []
                
                for n in map_data.get("nodes", []):
                    nodes.append(Node(id=str(n.get("id")), label=n.get("label", ""), title=n.get("title", ""), size=25))
                for e in map_data.get("edges", []):
                    edges.append(Edge(source=str(e.get("source")), target=str(e.get("target")), label=e.get("label", "")))
                    
                config = Config(width=600, height=400, directed=True, nodeHighlightBehavior=True, highlightColor="#F7A7A6", collapsible=True)
                agraph(nodes=nodes, edges=edges, config=config)
            except Exception as e:
                st.error(f"Error rendering mind map: {e}")
        
    with tab_report:
        st.subheader("Executive Report")
        if st.button("Generate Summary Report"):
            if not st.session_state.ingested_docs:
                st.warning("Please index documents first.")
            else:
                with st.spinner("Synthesizing report..."):
                    try:
                        results = retrieve("Summarize the main topics and key findings.", reranker_top_k=20)
                        st.session_state["report_md"] = generator.generate_report(results, st.session_state.backend)
                    except Exception as e:
                        st.error(f"Error: {e}")
        
        if "report_md" in st.session_state:
            st.markdown(st.session_state["report_md"])
            import markdown
            from xhtml2pdf import pisa
            import io
            html = markdown.markdown(st.session_state["report_md"])
            pdf_bytes = io.BytesIO()
            pisa.CreatePDF(io.StringIO(html), dest=pdf_bytes)
            st.download_button("📥 Download as PDF", data=pdf_bytes.getvalue(), file_name="Report.pdf", mime="application/pdf")

    with tab_flashcard:
        st.subheader("Key Terms")
        if st.button("Generate Flashcards"):
            if not st.session_state.ingested_docs:
                st.warning("Please index documents first.")
            else:
                with st.spinner("Extracting terminology..."):
                    try:
                        results = retrieve("Key terms and definitions", reranker_top_k=15)
                        import json
                        fc_json = generator.generate_flashcards(results, st.session_state.backend)
                        st.session_state["flashcards"] = json.loads(fc_json.replace('```json', '').replace('```', '').strip())
                    except Exception as e:
                        st.error(f"Error: {e}")
        
        if "flashcards" in st.session_state:
            for fc in st.session_state["flashcards"]:
                html_code = f'''
                <div class="flashcard">
                    <div class="flashcard-inner">
                        <div class="flashcard-front">{fc.get("front", "")}</div>
                        <div class="flashcard-back">{fc.get("back", "")}</div>
                    </div>
                </div>
                '''
                st.markdown(html_code, unsafe_allow_html=True)

    with tab_quiz:
        st.subheader("Knowledge Check")
        difficulty = st.selectbox("Difficulty:", ["Beginner", "Intermediate", "Advanced"])
        if st.button("Generate Quiz"):
            if not st.session_state.ingested_docs:
                st.warning("Please index documents first.")
            else:
                with st.spinner("Formulating questions..."):
                    try:
                        results = retrieve(f"Provide context to create a {difficulty} level quiz.", reranker_top_k=15)
                        import json
                        quiz_json = generator.generate_quiz(results, st.session_state.backend, difficulty)
                        st.session_state["quiz"] = json.loads(quiz_json.replace('```json', '').replace('```', '').strip())
                    except Exception as e:
                        st.error(f"Error: {e}")
                        
        if "quiz" in st.session_state:
            for i, q in enumerate(st.session_state["quiz"]):
                st.markdown(f"**Q{i+1}: {q.get('question','')}**")
                choice = st.radio("Select answer:", q.get('options',[]), key=f"quiz_q_{i}")
                if st.button(f"Check Answer", key=f"check_{i}"):
                    if choice == q.get("answer"):
                        st.success("Correct!")
                    else:
                        st.error(f"Incorrect. The correct answer was: {q.get('answer')}")
                    st.info(f"Explanation: {q.get('explanation', '')}")

    with tab_info:
        st.subheader("High-Quality Infographics")
        if st.button("Generate Infographic Concepts"):
            if not st.session_state.ingested_docs:
                st.warning("Please index documents first.")
            else:
                with st.spinner("Distilling concepts..."):
                    try:
                        results = retrieve("Summarize key concepts for infographics.", reranker_top_k=20)
                        import json
                        concepts_json = generator.generate_infographic_concepts(results, st.session_state.backend)
                        st.session_state["info_concepts"] = json.loads(concepts_json.replace('```json', '').replace('```', '').strip())
                        st.session_state["info_results"] = results
                    except Exception as e:
                        st.error(f"Error: {e}")
                        
        if "info_concepts" in st.session_state:
            selected_concept = st.radio("Select an Infographic Concept:", st.session_state["info_concepts"])
            if st.button("🎨 Render Infographic"):
                with st.spinner("Generating SVG..."):
                    try:
                        svg_code = generator.generate_svg(selected_concept, st.session_state["info_results"], st.session_state.backend)
                        st.session_state["current_svg"] = svg_code
                    except Exception as e:
                        st.error(f"Error: {e}")
                        
        if "current_svg" in st.session_state:
            svg = st.session_state["current_svg"]
            import base64
            b64 = base64.b64encode(svg.encode('utf-8')).decode("utf-8")
            html = f'<img src="data:image/svg+xml;base64,{b64}"/>'
            st.markdown(html, unsafe_allow_html=True)
            
            col_svg, col_png = st.columns(2)
            col_svg.download_button("📥 Download SVG", data=svg, file_name="infographic.svg", mime="image/svg+xml")
            
            try:
                import cairosvg
                png_data = cairosvg.svg2png(bytestring=svg.encode('utf-8'))
                col_png.download_button("📥 Download PNG", data=png_data, file_name="infographic.png", mime="image/png")
            except Exception as e:
                col_png.error(f"CairoSVG missing or error: {e}")
        
    with tab_data:
        st.subheader("Data Table")
        if st.button("Extract Data"):
            if not st.session_state.ingested_docs:
                st.warning("Please index documents first.")
            else:
                with st.spinner("Mining quantitative data..."):
                    try:
                        results = retrieve("Extract metrics, numbers, and quantitative data.", reranker_top_k=20)
                        st.session_state["data_table_md"] = generator.generate_data_table(results, st.session_state.backend)
                    except Exception as e:
                        st.error(f"Error: {e}")
                        
        if "data_table_md" in st.session_state:
            md_text = st.session_state["data_table_md"]
            lines = [l.strip() for l in md_text.split('\\n') if '|' in l]
            if len(lines) >= 3:
                import pandas as pd
                headers = [c.strip() for c in lines[0].split('|') if c.strip()]
                data = []
                for line in lines[2:]:
                    cells = [c.strip() for c in line.split('|') if c.strip()]
                    if len(cells) == len(headers):
                        data.append(cells)
                    elif cells:
                        data.append(cells + [""] * (len(headers) - len(cells)))
                if headers and data:
                    df = pd.DataFrame(data, columns=headers)
                    st.dataframe(df, use_container_width=True)
                else:
                    st.markdown(md_text)
            else:
                st.markdown(md_text)
