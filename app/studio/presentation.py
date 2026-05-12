import json
import re
import os
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Union, Dict, Any

from app.studio.engine import (
    CorpusProfile, 
    retrieve_all_for_studio, 
    analyze_corpus,
    format_chunks_for_generation, 
    call_llm
)
from app.retriever import retrieve, RetrievedChunk

@dataclass
class Slide:
    slide_number: int
    slide_type: str      # "title" | "agenda" | "insight" | "stat" | "quote" | "compare" | "process" | "timeline_item" | "section_break" | "summary"
    title: str
    content: Union[str, List[Any]]  # str for most types, list for bullets/compare/process/timeline
    visual_suggestion: str    # what graphic/chart to show
    speaker_notes: str        # full paragraph, not bullets
    citations: List[str]      # ["DocName, p.N"]
    emphasis: str             # key phrase to make LARGE

def parse_json(text: str) -> Any:
    """Utility to strip markdown and parse JSON."""
    try:
        clean = re.sub(r'```json\s*|\s*```', '', text).strip()
        return json.loads(clean)
    except Exception as e:
        print(f"JSON Parse Error: {e}")
        return None

def generate_presentation(
    topic: str = None,
    doc_names: List[str] = None,
    num_slides: int = 12,
    style: str = "executive",
    llm_backend: str = "gemini"
) -> List[Slide]:
    """
    Two-pass slide deck generation engine.
    Pass 1: Outline generation based on corpus analysis.
    Pass 2: Targeted content generation for each slide.
    """
    
    # STEP 1 — Retrieve + analyze
    all_chunks = retrieve_all_for_studio(doc_names, max_chunks=40)
    if not all_chunks:
        raise ValueError("No documents indexed. Please ingest PDFs first.")
    
    profile = analyze_corpus(all_chunks, llm_backend)

    # STEP 2 — First pass: generate slide OUTLINE only
    outline_system = "You are a senior consultant creating a professional presentation outline. Return JSON only."
    
    outline_prompt = f"""
Based on these source materials, create a {num_slides}-slide presentation outline on: 
{topic or "the key insights from these documents"}.

DOMAIN: {profile.domain}
KEY THEMES: {", ".join(profile.primary_topics)}
HAS DATA: {profile.has_numbers}
HAS TIMELINE: {profile.has_dates}

SLIDE TYPE RULES:
- Slide 1: Always "title" type
- Slide 2: "agenda" or "key_insight" hook
- Last slide: Always "summary" with 3 takeaways
- Vary types: mix stat, quote, compare, process, section_break
- NEVER have 3 consecutive "insight" (bullet) slides
- If has_numbers=true: include at least 2 "stat" slides
- If has_dates=true: include at least 1 "timeline_item" slide

Return a JSON array of outline items:
[
  {{
    "slide_number": 1,
    "slide_type": "title",
    "title": "<slide title>",
    "one_line_content": "<what this slide covers>"
  }},
  ...
]
Return JSON only.
"""

    outline_resp = call_llm(outline_system, outline_prompt, llm_backend)
    outline = parse_json(outline_resp)
    
    if not outline:
        # Simple fallback outline if LLM fails
        outline = [
            {"slide_number": 1, "slide_type": "title", "title": topic or "Analysis Report", "one_line_content": "Title slide"},
            {"slide_number": 2, "slide_type": "summary", "title": "Key Takeaways", "one_line_content": "Overview of findings"}
        ]

    # STEP 3 — Second pass: generate FULL CONTENT per slide
    final_slides = []
    content_system = "You are a professional presentation writer. Be punchy. One idea per slide. No fluff. Return JSON only."

    for item in outline:
        slide_num = item["slide_number"]
        slide_type = item["slide_type"]
        slide_title = item["title"]
        slide_purpose = item["one_line_content"]

        # Targeted retrieval for this specific slide
        search_query = f"{slide_title} {slide_purpose}"
        slide_chunks = retrieve(search_query, reranker_top_k=5)
        
        content_prompt = f"""
Generate full content for this slide:

SLIDE TYPE: {slide_type}
SLIDE TITLE: {slide_title}
SLIDE PURPOSE: {slide_purpose}
SLIDE NUMBER: {slide_num} of {len(outline)}

TYPE-SPECIFIC INSTRUCTIONS:
If "stat": 
  - content = one powerful number/percentage with unit
  - emphasis = that number (make it giant on slide)
  - include context sentence below the number
If "quote":
  - content = one memorable direct quote from sources
  - emphasis = key phrase within the quote to highlight
If "compare":
  - content = [{{ "label": "Option A", "points": ["...", "...", "..."] }}, {{ "label": "Option B", "points": ["...", "...", "..."] }}]
If "process":
  - content = ["Step 1: ...", "Step 2: ...", "Step 3: ..."] (Max 5)
If "insight" or "agenda":
  - content = ["• point 1", "• point 2", "• point 3"] (Max 3 bullets, max 12 words each)
If "timeline_item":
  - content = [{{ "date": "...", "event": "..." }}]
If "summary":
  - content = ["Takeaway 1: ...", "Takeaway 2: ...", "Takeaway 3: ..."]

SPEAKER NOTES RULES:
- Write 3-4 full sentences.
- Include what to say OUT LOUD.
- Explain "why this matters".
- End with transition to next slide.

SOURCE CHUNKS:
{format_chunks_for_generation(slide_chunks)}

Return JSON:
{{
  "slide_number": {slide_num},
  "slide_type": "{slide_type}",
  "title": "{slide_title}",
  "content": <str or list per type above>,
  "visual_suggestion": "<specific graphic/chart description>",
  "speaker_notes": "<full paragraph>",
  "citations": ["DocName, p.N"],
  "emphasis": "<key phrase>"
}}
"""
        slide_resp = call_llm(content_system, content_prompt, llm_backend)
        slide_data = parse_json(slide_resp)
        
        if slide_data:
            final_slides.append(Slide(
                slide_number=slide_data.get("slide_number", slide_num),
                slide_type=slide_data.get("slide_type", slide_type),
                title=slide_data.get("title", slide_title),
                content=slide_data.get("content", ""),
                visual_suggestion=slide_data.get("visual_suggestion", ""),
                speaker_notes=slide_data.get("speaker_notes", ""),
                citations=slide_data.get("citations", []),
                emphasis=slide_data.get("emphasis", "")
            ))

    return final_slides

def export_to_pptx(
    slides: List[Slide],
    output_path: Path
) -> Path:
    """
    Exports a list of Slide objects to a real .pptx file using python-pptx.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    
    # Colors
    TEAL = RGBColor(0, 128, 128)
    DARK_GRAY = RGBColor(64, 64, 64)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(240, 240, 240)

    for s in slides:
        # Select layout based on type
        if s.slide_type == "title":
            layout = prs.slide_layouts[0]
        elif s.slide_type == "compare":
            layout = prs.slide_layouts[4] # Comparison
        elif s.slide_type == "section_break":
            layout = prs.slide_layouts[2] # Section Header
        else:
            layout = prs.slide_layouts[1] # Title and Content
            
        slide = prs.slides.add_slide(layout)
        title_shape = slide.shapes.title
        
        # Style Title
        if title_shape:
            title_shape.text = s.title
            title_run = title_shape.text_frame.paragraphs[0].runs[0]
            title_run.font.name = 'DM Serif Display' # Fallback is usually Times New Roman or Garamond
            title_run.font.size = Pt(40)
            title_run.font.color.rgb = TEAL if s.slide_type != "section_break" else WHITE

        # Content Logic
        if s.slide_type == "title":
            # Subtitle for title slide
            subtitle = slide.placeholders[1]
            doc_names = ", ".join(list(set([c.split(",")[0] for c in s.citations])))
            subtitle.text = f"Generated from: {doc_names or 'Sources'}\nCreated by AskMyPDF Studio"
            
            # Teal accent bar
            left = 0
            top = prs.slide_height - Inches(0.2)
            width = prs.slide_width
            height = Inches(0.2)
            shape = slide.shapes.add_shape(1, left, top, width, height) # 1 is rectangle
            shape.fill.solid()
            shape.fill.fore_color.rgb = TEAL
            shape.line.color.rgb = TEAL

        elif s.slide_type == "stat":
            # Clear standard content placeholder if exists
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = ""
                
            # Giant Number
            left = Inches(1)
            top = Inches(2.5)
            width = prs.slide_width - Inches(2)
            height = Inches(2)
            txt_box = slide.shapes.add_textbox(left, top, width, height)
            tf = txt_box.text_frame
            tf.text = str(s.content)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(80)
            run.font.bold = True
            run.font.color.rgb = TEAL
            
            # Context Sentence
            p2 = tf.add_paragraph()
            p2.text = s.emphasis if s.emphasis and s.emphasis != s.content else ""
            p2.alignment = PP_ALIGN.CENTER
            p2.font.size = Pt(18)
            p2.font.color.rgb = DARK_GRAY

        elif s.slide_type == "quote":
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = ""
            
            left = Inches(1.5)
            top = Inches(2)
            width = prs.slide_width - Inches(3)
            height = Inches(3)
            txt_box = slide.shapes.add_textbox(left, top, width, height)
            tf = txt_box.text_frame
            tf.word_wrap = True
            
            # Opening Quote mark
            p_mark = tf.paragraphs[0]
            p_mark.text = '"'
            p_mark.font.size = Pt(60)
            p_mark.font.color.rgb = TEAL
            
            # Quote Text
            p_quote = tf.add_paragraph()
            p_quote.text = str(s.content)
            p_quote.font.italic = True
            p_quote.font.size = Pt(24)
            
            # Source
            p_src = tf.add_paragraph()
            p_src.alignment = PP_ALIGN.RIGHT
            p_src.text = f"— {s.citations[0]}" if s.citations else ""
            p_src.font.size = Pt(14)

        elif s.slide_type == "compare":
            # content: [{"label": "A", "points": [...]}, {"label": "B", "points": [...]}]
            if isinstance(s.content, list) and len(s.content) >= 2:
                # Layout 4 has two placeholders for content
                left_col = slide.placeholders[1]
                right_col = slide.placeholders[2]
                
                left_col.text = s.content[0].get("label", "Option A")
                tf_l = left_col.text_frame
                for p_text in s.content[0].get("points", []):
                    p = tf_l.add_paragraph()
                    p.text = p_text
                    p.level = 0
                
                right_col.text = s.content[1].get("label", "Option B")
                tf_r = right_col.text_frame
                for p_text in s.content[1].get("points", []):
                    p = tf_r.add_paragraph()
                    p.text = p_text
                    p.level = 0

        elif s.slide_type == "timeline_item":
            # content: [{"date": "...", "event": "..."}]
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            if isinstance(s.content, list):
                for item in s.content:
                    p = tf.add_paragraph()
                    p.text = f"{item.get('date', '')}: {item.get('event', '')}"
                    p.level = 0
                    p.font.bold = True
                    p.font.color.rgb = TEAL

        elif s.slide_type == "summary":
            # 3 takeaways in teal circles
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "" # Clear placeholder
                
            takeaways = s.content if isinstance(s.content, list) else [s.content]
            for i, text in enumerate(takeaways[:3]):
                y_pos = Inches(2 + i * 1.5)
                # Circle
                shape = slide.shapes.add_shape(9, Inches(1), y_pos, Inches(0.8), Inches(0.8)) # 9 is Oval
                shape.fill.solid()
                shape.fill.fore_color.rgb = TEAL
                shape.line.visible = False
                tf = shape.text_frame
                tf.text = str(i + 1)
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                p.runs[0].font.size = Pt(24)
                p.runs[0].font.color.rgb = WHITE
                
                # Text
                txt = slide.shapes.add_textbox(Inches(2), y_pos + Inches(0.1), Inches(7), Inches(0.8))
                tf2 = txt.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = str(text)
                p2.font.size = Pt(22)
                p2.font.color.rgb = DARK_GRAY

        elif s.slide_type == "section_break":
            # Dark background
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = DARK_GRAY
            # Accent line
            left = prs.slide_width * 0.1
            top = prs.slide_height * 0.6
            width = prs.slide_width * 0.8
            height = Pt(4)
            line = slide.shapes.add_shape(1, left, top, width, height)
            line.fill.solid()
            line.fill.fore_color.rgb = TEAL
            line.line.visible = False

        else: # insight, agenda, summary, process, timeline
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            if isinstance(s.content, list):
                for i, item in enumerate(s.content):
                    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                    p.text = str(item)
                    p.level = 0
            else:
                body_shape.text = str(s.content)

        # Add Speaker Notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = s.speaker_notes

        # Footer (added manually as shapes)
        footer_y = prs.slide_height - Inches(0.4)
        footer_text = f"{s.slide_number} | {s.title} | AskMyPDF Studio"
        txBox = slide.shapes.add_textbox(Inches(0.5), footer_y, prs.slide_width - Inches(1), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_GRAY

    prs.save(str(output_path))
    return output_path
